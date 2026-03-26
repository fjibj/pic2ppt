"""
基础形状处理器
包含：rect, circle, ellipse, line
"""

from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE

from . import ElementHandler, RenderContext
from ..models import SVGElement, ShapeType
from ..color_utils import ColorParser


class BaseHandler(ElementHandler):
    """处理器基类，提供通用工具方法"""

    def _apply_style(self, shape, element: SVGElement, context: RenderContext):
        """应用样式到形状"""
        style = context.parent_style.merge(element.style)

        # 填充 - 处理明确填充、无填充和默认填充的情况
        if style.fill and style.fill.lower() == 'none':
            # 明确指定无填充
            shape.fill.background()
        elif style.fill:
            # 有明确填充颜色
            fill_color = ColorParser.parse(style.fill)
            if fill_color:
                shape.fill.solid()
                shape.fill.fore_color.rgb = fill_color
        else:
            # 默认填充（SVG默认是黑色填充，但通常图表元素有明确填充）
            # 如果元素是rect/circle/ellipse且有stroke但没有fill，使用白色填充
            if element.tag in ('rect', 'circle', 'ellipse', 'polygon'):
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(255, 255, 255)  # 默认白色填充
            else:
                shape.fill.background()

        # 边框
        if style.stroke and style.stroke.lower() != 'none':
            stroke_color = ColorParser.parse(style.stroke)
            if stroke_color:
                shape.line.color.rgb = stroke_color
                shape.line.width = Pt(ColorParser.parse_length(str(style.stroke_width)) or 1)

            # 虚线
            if style.stroke_dasharray:
                shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        else:
            shape.line.fill.background()


class RectHandler(BaseHandler):
    """矩形处理器"""

    def can_handle(self, element: SVGElement) -> bool:
        return element.tag == 'rect'

    def handle(self, element: SVGElement, context: RenderContext):
        x = element.get_float_attr('x')
        y = element.get_float_attr('y')
        width = element.get_float_attr('width')
        height = element.get_float_attr('height')
        rx = element.get_float_attr('rx')
        ry = element.get_float_attr('ry', rx)  # 如果没有ry，使用rx

        # 计算PPT坐标
        left = Inches(context.offset_x + context.svg_to_inches(x))
        top = Inches(context.offset_y + context.svg_to_inches(y))
        w = Inches(context.svg_to_inches(width))
        h = Inches(context.svg_to_inches(height))

        # 选择形状类型
        if rx > 0 or ry > 0:
            shape = context.slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
        else:
            shape = context.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)

        self._apply_style(shape, element, context)
        return shape


class CircleHandler(BaseHandler):
    """圆形处理器"""

    def can_handle(self, element: SVGElement) -> bool:
        return element.tag == 'circle'

    def handle(self, element: SVGElement, context: RenderContext):
        cx = element.get_float_attr('cx')
        cy = element.get_float_attr('cy')
        r = element.get_float_attr('r')

        # 圆的包围盒
        left = Inches(context.offset_x + context.svg_to_inches(cx - r))
        top = Inches(context.offset_y + context.svg_to_inches(cy - r))
        size = Inches(context.svg_to_inches(r * 2))

        shape = context.slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
        self._apply_style(shape, element, context)
        return shape


class EllipseHandler(BaseHandler):
    """椭圆处理器"""

    def can_handle(self, element: SVGElement) -> bool:
        return element.tag == 'ellipse'

    def handle(self, element: SVGElement, context: RenderContext):
        cx = element.get_float_attr('cx')
        cy = element.get_float_attr('cy')
        rx = element.get_float_attr('rx')
        ry = element.get_float_attr('ry')

        left = Inches(context.offset_x + context.svg_to_inches(cx - rx))
        top = Inches(context.offset_y + context.svg_to_inches(cy - ry))
        width = Inches(context.svg_to_inches(rx * 2))
        height = Inches(context.svg_to_inches(ry * 2))

        shape = context.slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
        self._apply_style(shape, element, context)
        return shape


class LineHandler(BaseHandler):
    """直线处理器 - 支持内置箭头样式"""

    def can_handle(self, element: SVGElement) -> bool:
        return element.tag == 'line'

    def handle(self, element: SVGElement, context: RenderContext):
        x1 = element.get_float_attr('x1')
        y1 = element.get_float_attr('y1')
        x2 = element.get_float_attr('x2')
        y2 = element.get_float_attr('y2')

        begin_x = Inches(context.offset_x + context.svg_to_inches(x1))
        begin_y = Inches(context.offset_y + context.svg_to_inches(y1))
        end_x = Inches(context.offset_x + context.svg_to_inches(x2))
        end_y = Inches(context.offset_y + context.svg_to_inches(y2))

        # 检查是否有箭头
        marker_end = element.attrib.get('marker-end', '')
        marker_start = element.attrib.get('marker-start', '')
        has_arrow_end = bool(marker_end) and ('arrow' in marker_end.lower() or 'url(#' in marker_end)
        has_arrow_start = bool(marker_start) and ('arrow' in marker_start.lower() or 'url(#' in marker_start)

        # 创建线条
        connector = context.slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, begin_x, begin_y, end_x, end_y
        )

        # 应用样式（包括箭头）
        style = element.style.merge(context.parent_style)
        self._apply_connector_style(connector, style, has_arrow_start, has_arrow_end)

        return connector

    def _apply_connector_style(self, connector, style, has_arrow_start: bool, has_arrow_end: bool):
        """应用样式到连接器，包括箭头"""
        from lxml import etree

        # 线条颜色
        if style.stroke and style.stroke.lower() != 'none':
            stroke_color = ColorParser.parse(style.stroke)
            if stroke_color:
                connector.line.color.rgb = stroke_color
        else:
            connector.line.color.rgb = RGBColor(0, 0, 0)

        # 线条宽度
        line_width = Pt(ColorParser.parse_length(str(style.stroke_width)) or 1)
        connector.line.width = line_width

        # 虚线样式
        if style.stroke_dasharray:
            connector.line.dash_style = MSO_LINE_DASH_STYLE.DASH

        # 箭头样式 - 使用XML添加
        if has_arrow_start or has_arrow_end:
            P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
            A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'

            spPr = connector.element.find('{%s}spPr' % P_NS)
            if spPr is not None:
                ln = spPr.find('{%s}ln' % A_NS)
                if ln is not None:
                    # 根据线条宽度选择箭头大小
                    # 粗线使用大箭头，细线使用中箭头
                    if line_width >= Pt(3):
                        arrow_size = 'lg'  # 大号箭头
                    elif line_width >= Pt(1.5):
                        arrow_size = 'med'  # 中号箭头
                    else:
                        arrow_size = 'sm'  # 小号箭头

                    # PPT中：headEnd = 起点箭头, tailEnd = 终点箭头
                    # 注意：这与直觉相反！
                    if has_arrow_end:
                        tailEnd = etree.SubElement(ln, '{%s}tailEnd' % A_NS)
                        tailEnd.set('type', 'arrow')
                        tailEnd.set('w', arrow_size)
                        tailEnd.set('len', arrow_size)
                    if has_arrow_start:
                        headEnd = etree.SubElement(ln, '{%s}headEnd' % A_NS)
                        headEnd.set('type', 'arrow')
                        headEnd.set('w', arrow_size)
                        headEnd.set('len', arrow_size)
