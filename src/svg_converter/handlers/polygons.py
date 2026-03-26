"""
多边形处理器
包含：polygon, polyline, path
"""

import math
from typing import Optional, Tuple, List
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE

from . import ElementHandler, RenderContext
from ..models import SVGElement, ShapeType
from ..geometry import GeometryAnalyzer
from ..color_utils import ColorParser
from .basic_shapes import BaseHandler


class PolygonHandler(BaseHandler):
    """多边形处理器 - 使用几何分析自动识别形状"""

    # 形状类型到PPT形状的映射
    SHAPE_MAP = {
        ShapeType.RECTANGLE: MSO_SHAPE.RECTANGLE,
        ShapeType.ROUNDED_RECTANGLE: MSO_SHAPE.ROUNDED_RECTANGLE,
        ShapeType.OVAL: MSO_SHAPE.OVAL,
        ShapeType.DIAMOND: MSO_SHAPE.DIAMOND,
        ShapeType.ISOSCELES_TRIANGLE: MSO_SHAPE.ISOSCELES_TRIANGLE,
        ShapeType.RIGHT_TRIANGLE: MSO_SHAPE.RIGHT_TRIANGLE,
        ShapeType.HEXAGON: MSO_SHAPE.HEXAGON,
        ShapeType.OCTAGON: MSO_SHAPE.OCTAGON,
        ShapeType.PENTAGON: MSO_SHAPE.REGULAR_PENTAGON,
        ShapeType.TRAPEZOID: MSO_SHAPE.TRAPEZOID,
        ShapeType.PARALLELOGRAM: MSO_SHAPE.PARALLELOGRAM,
    }

    def can_handle(self, element: SVGElement) -> bool:
        return element.tag in ('polygon', 'polyline')

    def handle(self, element: SVGElement, context: RenderContext):
        points_str = element.attrib.get('points', '')
        analyzer = GeometryAnalyzer()
        points = analyzer.parse_points(points_str)

        if not points:
            return None

        # 区分 polygon 和 polyline
        if element.tag == 'polyline':
            return self._handle_polyline(points, element, context)
        else:
            return self._handle_polygon(points, element, context)

    def _handle_polygon(self, points: List[Tuple[float, float]], element: SVGElement, context: RenderContext):
        """处理多边形（闭合形状）"""
        analyzer = GeometryAnalyzer()
        shape_type, bbox = analyzer.analyze_points(points)

        # 转换为PPT坐标
        left = Inches(context.offset_x + context.svg_to_inches(bbox.x))
        top = Inches(context.offset_y + context.svg_to_inches(bbox.y))
        width = Inches(context.svg_to_inches(bbox.width))
        height = Inches(context.svg_to_inches(bbox.height))

        # 创建形状
        if shape_type in self.SHAPE_MAP:
            shape = context.slide.shapes.add_shape(
                self.SHAPE_MAP[shape_type], left, top, width, height
            )
            self._apply_style(shape, element, context)
            return shape
        elif shape_type == ShapeType.STRAIGHT_LINE and len(points) == 2:
            # 简单线条
            x1, y1 = points[0]
            x2, y2 = points[1]
            begin_x = Inches(context.offset_x + context.svg_to_inches(x1))
            begin_y = Inches(context.offset_y + context.svg_to_inches(y1))
            end_x = Inches(context.offset_x + context.svg_to_inches(x2))
            end_y = Inches(context.offset_y + context.svg_to_inches(y2))

            connector = context.slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT, begin_x, begin_y, end_x, end_y
            )
            self._apply_style(connector, element, context)
            return connector
        else:
            # 复杂多边形 - 使用矩形代替
            shape = context.slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, left, top, width, height
            )
            self._apply_style(shape, element, context)
            return shape

    def _handle_polyline(self, points: List[Tuple[float, float]], element: SVGElement, context: RenderContext):
        """处理折线（不闭合的多段线）"""
        if len(points) < 2:
            return None

        # 获取样式
        style = context.parent_style.merge(element.style)

        # 检查是否有箭头
        marker_end = element.attrib.get('marker-end', '')
        marker_start = element.attrib.get('marker-start', '')
        has_arrow_end = bool(marker_end) and ('arrow' in marker_end.lower() or 'url(#' in marker_end)
        has_arrow_start = bool(marker_start) and ('arrow' in marker_start.lower() or 'url(#' in marker_start)

        if len(points) == 2:
            # 两点折线 = 直线
            x1, y1 = points[0]
            x2, y2 = points[1]
            begin_x = Inches(context.offset_x + context.svg_to_inches(x1))
            begin_y = Inches(context.offset_y + context.svg_to_inches(y1))
            end_x = Inches(context.offset_x + context.svg_to_inches(x2))
            end_y = Inches(context.offset_y + context.svg_to_inches(y2))

            connector = context.slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT, begin_x, begin_y, end_x, end_y
            )

            # 应用样式
            self._apply_line_style(connector, style)

            return connector
        else:
            # 多点折线 - 创建为自由形状（freeform）来保持折线形状
            return self._create_freeform_polyline(points, element, context, has_arrow_start, has_arrow_end)

    def _create_freeform_polyline(self, points: List[Tuple[float, float]], element: SVGElement,
                                   context: RenderContext, has_arrow_start: bool, has_arrow_end: bool):
        """创建折线 - 使用直线连接器拼接（避免FreeformBuilder的bug）"""
        from pptx.enum.shapes import MSO_CONNECTOR

        if not points or len(points) < 2:
            return None

        style = context.parent_style.merge(element.style)

        shapes = []

        # 为每条线段创建连接器
        for i in range(len(points) - 1):
            x1, y1 = points[i]
            x2, y2 = points[i + 1]

            begin_x = Inches(context.offset_x + context.svg_to_inches(x1))
            begin_y = Inches(context.offset_y + context.svg_to_inches(y1))
            end_x = Inches(context.offset_x + context.svg_to_inches(x2))
            end_y = Inches(context.offset_y + context.svg_to_inches(y2))

            # 创建直线连接器
            connector = context.slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT, begin_x, begin_y, end_x, end_y
            )

            # 应用线条样式
            self._apply_line_style(connector, style)

            shapes.append(connector)

        # 返回第一个形状作为主形状（后续改进：返回组）
        return shapes[0] if shapes else None

    def _apply_line_style(self, connector, style):
        """应用线条样式到连接器"""
        if style.stroke and style.stroke.lower() != 'none':
            stroke_color = ColorParser.parse(style.stroke)
            if stroke_color:
                connector.line.color.rgb = stroke_color
        else:
            connector.line.color.rgb = RGBColor(0, 0, 0)

        line_width = Pt(ColorParser.parse_length(str(style.stroke_width)) or 1)
        connector.line.width = line_width

        if style.stroke_dasharray:
            connector.line.dash_style = MSO_LINE_DASH_STYLE.DASH

    def _apply_line_style_to_shape(self, shape, style):
        """应用线条样式到形状（用于自由形状）"""
        if style.stroke and style.stroke.lower() != 'none':
            stroke_color = ColorParser.parse(style.stroke)
            if stroke_color:
                shape.line.color.rgb = stroke_color
        else:
            shape.line.color.rgb = RGBColor(0, 0, 0)

        shape.line.width = Pt(ColorParser.parse_length(str(style.stroke_width)) or 1)

        if style.stroke_dasharray:
            shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH

        # 填充处理 - 与基本形状一致
        if style.fill and style.fill.lower() == 'none':
            shape.fill.background()
        elif style.fill:
            fill_color = ColorParser.parse(style.fill)
            if fill_color:
                shape.fill.solid()
                shape.fill.fore_color.rgb = fill_color
        else:
            # 默认白色填充
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(255, 255, 255)


class PathHandler(BaseHandler):
    """路径处理器 - 支持直线、曲线、矩形、箭头等"""

    def can_handle(self, element: SVGElement) -> bool:
        return element.tag == 'path'

    def handle(self, element: SVGElement, context: RenderContext):
        d = element.attrib.get('d', '')
        if not d:
            return None

        # 获取样式和箭头信息
        style = context.parent_style.merge(element.style)
        marker_end = element.attrib.get('marker-end', '')
        marker_start = element.attrib.get('marker-start', '')
        has_arrow_end = bool(marker_end) and ('arrow' in marker_end.lower() or 'url(#' in marker_end)
        has_arrow_start = bool(marker_start) and ('arrow' in marker_start.lower() or 'url(#' in marker_start)

        # 首先尝试解析为简单直线（M ... L ...）
        line_points = self._parse_line_path(d)
        if line_points:
            x1, y1, x2, y2 = line_points
            begin_x = Inches(context.offset_x + context.svg_to_inches(x1))
            begin_y = Inches(context.offset_y + context.svg_to_inches(y1))
            end_x = Inches(context.offset_x + context.svg_to_inches(x2))
            end_y = Inches(context.offset_y + context.svg_to_inches(y2))

            # 创建线条
            connector = context.slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT, begin_x, begin_y, end_x, end_y
            )

            # 应用样式
            self._apply_line_style(connector, style)

            return connector

        # 尝试解析为曲线路径
        curve_result = self._parse_curve_path(d)
        if curve_result:
            return self._handle_curve_path(curve_result, element, context, has_arrow_start, has_arrow_end)

        # 尝试解析为矩形路径
        if self._is_rect_path(d):
            analyzer = GeometryAnalyzer()
            shape_type, bbox = analyzer.analyze_path(d)

            left = Inches(context.offset_x + context.svg_to_inches(bbox.x))
            top = Inches(context.offset_y + context.svg_to_inches(bbox.y))
            width = Inches(context.svg_to_inches(bbox.width))
            height = Inches(context.svg_to_inches(bbox.height))

            shape = context.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
            self._apply_style(shape, element, context)
            return shape

        # 复杂路径 - 尝试提取所有点并创建自由形状
        return self._create_freeform_from_path(d, element, context)

    def _apply_line_style(self, connector, style):
        """应用线条样式到连接器"""
        if style.stroke and style.stroke.lower() != 'none':
            stroke_color = ColorParser.parse(style.stroke)
            if stroke_color:
                connector.line.color.rgb = stroke_color
        else:
            connector.line.color.rgb = RGBColor(0x7f, 0x8c, 0x8d)

        connector.line.width = Pt(ColorParser.parse_length(str(style.stroke_width)) or 1.5)

        if style.stroke_dasharray:
            connector.line.dash_style = MSO_LINE_DASH_STYLE.DASH

    def _parse_line_path(self, d: str) -> Optional[Tuple[float, float, float, float]]:
        """
        解析简单的直线路径，如 M 600 100 L 600 130
        返回 (x1, y1, x2, y2) 或 None
        """
        import re

        # 简化路径，移除多余空格
        d = d.strip()

        # 匹配 M x y L x y 模式（支持各种空格和分隔符）
        line_pattern = r'^[Mm]\s*(-?[\d.]+)[,\s]+(-?[\d.]+)[,\s]*[Ll]\s*(-?[\d.]+)[,\s]+(-?[\d.]+)$'
        match = re.match(line_pattern, d)

        if match:
            x1, y1, x2, y2 = map(float, match.groups())
            return (x1, y1, x2, y2)

        return None

    def _parse_curve_path(self, d: str) -> Optional[dict]:
        """
        解析曲线路径，支持 Q(二次贝塞尔)、C(三次贝塞尔)、A(圆弧)
        返回路径信息字典或None
        """
        import re

        d = d.strip()

        # 检查是否包含曲线命令
        has_quadratic = 'Q' in d or 'q' in d  # 二次贝塞尔
        has_cubic = 'C' in d or 'c' in d      # 三次贝塞尔
        has_arc = 'A' in d or 'a' in d        # 圆弧

        if not (has_quadratic or has_cubic or has_arc):
            return None

        # 提取所有命令和参数
        commands = re.findall(r'([MmLlHhVvCcSsQqTtAaZz])\s*([^MmLlHhVvCcSsQqTtAaZz]*)', d)

        if not commands:
            return None

        return {
            'type': 'curve',
            'has_quadratic': has_quadratic,
            'has_cubic': has_cubic,
            'has_arc': has_arc,
            'commands': commands,
            'original': d
        }

    def _handle_curve_path(self, curve_info: dict, element: SVGElement, context: RenderContext,
                           has_arrow_start: bool, has_arrow_end: bool):
        """处理曲线路径 - 使用Curve连接器（如果支持）或线段逼近"""
        from pptx.enum.shapes import MSO_CONNECTOR
        import re

        d = curve_info['original']
        style = context.parent_style.merge(element.style)

        # 首先尝试解析为简单的二次贝塞尔曲线 (M x y Q cx cy x y)
        quad_match = re.match(r'^[Mm]\s*(-?[\d.]+)[,\s]+(-?[\d.]+)[,\s]*[Qq]\s*(-?[\d.]+)[,\s]+(-?[\d.]+)[,\s]*(-?[\d.]+)[,\s]+(-?[\d.]+)$', d.strip())
        if quad_match:
            # 简单的二次贝塞尔曲线
            x1, y1, cx, cy, x2, y2 = map(float, quad_match.groups())
            # 使用多点线段逼近曲线（更多点以获得更平滑的效果）
            points = self._subdivide_quadratic_bezier((x1, y1), (cx, cy), (x2, y2), steps=25)
            return self._create_curve_from_points(points, element, context, has_arrow_start, has_arrow_end)

        # 尝试解析为三次贝塞尔曲线
        cubic_match = re.match(r'^[Mm]\s*(-?[\d.]+)[,\s]+(-?[\d.]+)[,\s]*[Cc]\s*(-?[\d.]+)[,\s]+(-?[\d.]+)[,\s]*(-?[\d.]+)[,\s]+(-?[\d.]+)[,\s]*(-?[\d.]+)[,\s]+(-?[\d.]+)$', d.strip())
        if cubic_match:
            x1, y1, cx1, cy1, cx2, cy2, x2, y2 = map(float, cubic_match.groups())
            points = self._subdivide_cubic_bezier((x1, y1), (cx1, cy1), (cx2, cy2), (x2, y2), steps=25)
            return self._create_curve_from_points(points, element, context, has_arrow_start, has_arrow_end)

        # 通用曲线处理
        if curve_info.get('has_arc'):
            points = self._parse_curve_to_points(d, steps=25)
        else:
            points = self._parse_curve_to_points(d, steps=20)

        if len(points) < 2:
            return None

        return self._create_curve_from_points(points, element, context, has_arrow_start, has_arrow_end)

    def _create_curve_from_points(self, points: List[Tuple[float, float]], element: SVGElement,
                                   context: RenderContext, has_arrow_start: bool, has_arrow_end: bool):
        """从点序列创建曲线 - 使用线条拼接来模拟曲线"""
        from pptx.enum.shapes import MSO_CONNECTOR
        from lxml import etree

        if len(points) < 2:
            return None

        style = element.style.merge(context.parent_style)
        first_shape = None
        last_shape = None

        # 为每条线段创建连接器
        for i in range(len(points) - 1):
            x1, y1 = points[i]
            x2, y2 = points[i + 1]

            begin_x = Inches(context.offset_x + context.svg_to_inches(x1))
            begin_y = Inches(context.offset_y + context.svg_to_inches(y1))
            end_x = Inches(context.offset_x + context.svg_to_inches(x2))
            end_y = Inches(context.offset_y + context.svg_to_inches(y2))

            # 创建直线连接器
            connector = context.slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT, begin_x, begin_y, end_x, end_y
            )

            # 应用样式
            self._apply_line_style(connector, style)

            # 只在第一个线段添加起始箭头，最后一个线段添加结束箭头
            if i == 0 and has_arrow_start:
                self._add_arrow_to_connector(connector, 'tail', style.stroke_width)
            if i == len(points) - 2 and has_arrow_end:
                self._add_arrow_to_connector(connector, 'head', style.stroke_width)

            if first_shape is None:
                first_shape = connector
            last_shape = connector

        return first_shape

    def _add_arrow_to_connector(self, connector, end_type: str, line_width=None):
        """为连接器添加箭头"""
        from lxml import etree
        from pptx.util import Pt
        P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'

        # 根据线条宽度选择箭头大小
        if line_width is None:
            line_width = connector.line.width

        if line_width >= Pt(3):
            arrow_size = 'lg'
        elif line_width >= Pt(1.5):
            arrow_size = 'med'
        else:
            arrow_size = 'sm'

        spPr = connector.element.find('{%s}spPr' % P_NS)
        if spPr is not None:
            ln = spPr.find('{%s}ln' % A_NS)
            if ln is not None:
                # PPT中：headEnd = 起点箭头, tailEnd = 终点箭头
                if end_type == 'head':  # 终点箭头 → tailEnd
                    tailEnd = etree.SubElement(ln, '{%s}tailEnd' % A_NS)
                    tailEnd.set('type', 'arrow')
                    tailEnd.set('w', arrow_size)
                    tailEnd.set('len', arrow_size)
                else:  # 起点箭头 → headEnd
                    headEnd = etree.SubElement(ln, '{%s}headEnd' % A_NS)
                    headEnd.set('type', 'arrow')
                    headEnd.set('w', arrow_size)
                    headEnd.set('len', arrow_size)

    def _parse_curve_to_points(self, d: str, steps: int = 10) -> List[Tuple[float, float]]:
        """解析SVG路径，将曲线细分为密集的点序列"""
        import re

        points = []
        current_pos = (0, 0)
        start_pos = (0, 0)
        last_cp = None  # 上次使用的控制点（用于平滑曲线）

        # 解析所有命令
        commands = re.findall(r'([MmLlHhVvCcSsQqTtAaZz])\s*([^MmLlHhVvCcSsQqTtAaZz]*)', d)

        for cmd, params_str in commands:
            params = self._parse_numbers(params_str)

            if cmd == 'M' or cmd == 'm':  # 移动到
                if len(params) >= 2:
                    if cmd == 'M':
                        current_pos = (params[0], params[1])
                    else:  # 相对坐标
                        current_pos = (current_pos[0] + params[0], current_pos[1] + params[1])
                    start_pos = current_pos
                    points.append(current_pos)
                    # 后续参数视为L命令
                    for i in range(2, len(params), 2):
                        if i + 1 < len(params):
                            current_pos = (params[i], params[i + 1])
                            points.append(current_pos)

            elif cmd == 'L' or cmd == 'l':  # 直线到
                for i in range(0, len(params), 2):
                    if i + 1 < len(params):
                        if cmd == 'L':
                            current_pos = (params[i], params[i + 1])
                        else:
                            current_pos = (current_pos[0] + params[i], current_pos[1] + params[i + 1])
                        points.append(current_pos)

            elif cmd == 'H' or cmd == 'h':  # 水平线
                for x in params:
                    if cmd == 'H':
                        current_pos = (x, current_pos[1])
                    else:
                        current_pos = (current_pos[0] + x, current_pos[1])
                    points.append(current_pos)

            elif cmd == 'V' or cmd == 'v':  # 垂直线
                for y in params:
                    if cmd == 'V':
                        current_pos = (current_pos[0], y)
                    else:
                        current_pos = (current_pos[0], current_pos[1] + y)
                    points.append(current_pos)

            elif cmd == 'Q' or cmd == 'q':  # 二次贝塞尔曲线
                for i in range(0, len(params), 4):
                    if i + 3 < len(params):
                        if cmd == 'Q':
                            cp = (params[i], params[i + 1])
                            end = (params[i + 2], params[i + 3])
                        else:
                            cp = (current_pos[0] + params[i], current_pos[1] + params[i + 1])
                            end = (current_pos[0] + params[i + 2], current_pos[1] + params[i + 3])
                        # 细分曲线为多个点
                        curve_points = self._subdivide_quadratic_bezier(current_pos, cp, end, steps)
                        points.extend(curve_points[1:])  # 跳过第一个点（已经是current_pos）
                        current_pos = end
                        last_cp = cp

            elif cmd == 'T' or cmd == 't':  # 平滑二次贝塞尔曲线
                for i in range(0, len(params), 2):
                    if i + 1 < len(params):
                        if cmd == 'T':
                            end = (params[i], params[i + 1])
                        else:
                            end = (current_pos[0] + params[i], current_pos[1] + params[i + 1])
                        # 计算控制点（反射上一个控制点）
                        if last_cp:
                            cp = (2 * current_pos[0] - last_cp[0], 2 * current_pos[1] - last_cp[1])
                        else:
                            cp = current_pos
                        curve_points = self._subdivide_quadratic_bezier(current_pos, cp, end, steps)
                        points.extend(curve_points[1:])
                        current_pos = end
                        last_cp = cp

            elif cmd == 'C' or cmd == 'c':  # 三次贝塞尔曲线
                for i in range(0, len(params), 6):
                    if i + 5 < len(params):
                        if cmd == 'C':
                            cp1 = (params[i], params[i + 1])
                            cp2 = (params[i + 2], params[i + 3])
                            end = (params[i + 4], params[i + 5])
                        else:
                            cp1 = (current_pos[0] + params[i], current_pos[1] + params[i + 1])
                            cp2 = (current_pos[0] + params[i + 2], current_pos[1] + params[i + 3])
                            end = (current_pos[0] + params[i + 4], current_pos[1] + params[i + 5])
                        curve_points = self._subdivide_cubic_bezier(current_pos, cp1, cp2, end, steps)
                        points.extend(curve_points[1:])
                        current_pos = end
                        last_cp = cp2

            elif cmd == 'A' or cmd == 'a':  # 圆弧
                for i in range(0, len(params), 7):
                    if i + 6 < len(params):
                        rx, ry, x_rot, large_arc, sweep, x, y = params[i:i + 7]
                        if cmd == 'A':
                            end = (x, y)
                        else:
                            end = (current_pos[0] + x, current_pos[1] + y)
                        large_arc = int(large_arc)
                        sweep = int(sweep)
                        arc_points = self._subdivide_arc(current_pos, rx, ry, x_rot, large_arc, sweep, end, steps)
                        points.extend(arc_points[1:])
                        current_pos = end

            elif cmd == 'Z' or cmd == 'z':  # 闭合路径
                if points and start_pos != current_pos:
                    points.append(start_pos)
                current_pos = start_pos

        return points

    def _parse_numbers(self, s: str) -> List[float]:
        """解析字符串中的数字"""
        import re
        numbers = re.findall(r'-?\d+\.?\d*', s)
        return [float(n) for n in numbers]

    def _subdivide_quadratic_bezier(self, p0, p1, p2, steps):
        """细分二次贝塞尔曲线为直线段"""
        points = []
        for i in range(steps + 1):
            t = i / steps
            # 二次贝塞尔公式: B(t) = (1-t)^2 * P0 + 2(1-t)t * P1 + t^2 * P2
            x = (1 - t) ** 2 * p0[0] + 2 * (1 - t) * t * p1[0] + t ** 2 * p2[0]
            y = (1 - t) ** 2 * p0[1] + 2 * (1 - t) * t * p1[1] + t ** 2 * p2[1]
            points.append((x, y))
        return points

    def _subdivide_cubic_bezier(self, p0, p1, p2, p3, steps):
        """细分三次贝塞尔曲线为直线段"""
        points = []
        for i in range(steps + 1):
            t = i / steps
            # 三次贝塞尔公式
            mt = 1 - t
            x = mt ** 3 * p0[0] + 3 * mt ** 2 * t * p1[0] + 3 * mt * t ** 2 * p2[0] + t ** 3 * p3[0]
            y = mt ** 3 * p0[1] + 3 * mt ** 2 * t * p1[1] + 3 * mt * t ** 2 * p2[1] + t ** 3 * p3[1]
            points.append((x, y))
        return points

    def _subdivide_arc(self, start, rx, ry, x_rot, large_arc, sweep, end, steps):
        """细分椭圆弧为直线段 - 正确实现SVG椭圆弧参数化"""
        import math

        # 如果rx或ry为0，退化为直线
        if rx == 0 or ry == 0:
            return [start, end]

        # 确保rx, ry为正
        rx = abs(rx)
        ry = abs(ry)

        # 1. 将坐标转换到椭圆坐标系
        x1, y1 = start
        x2, y2 = end

        # 转换为弧度
        phi = math.radians(x_rot)

        # 2. 计算中点差值
        dx2 = (x1 - x2) / 2
        dy2 = (y1 - y2) / 2

        # 3. 旋转到椭圆坐标系
        x1p = math.cos(phi) * dx2 + math.sin(phi) * dy2
        y1p = -math.sin(phi) * dx2 + math.cos(phi) * dy2

        # 4. 确保半径足够大
        lambda_val = (x1p ** 2) / (rx ** 2) + (y1p ** 2) / (ry ** 2)
        if lambda_val > 1:
            sqrt_lambda = math.sqrt(lambda_val)
            rx *= sqrt_lambda
            ry *= sqrt_lambda

        # 5. 计算中心点
        sign = -1 if large_arc == sweep else 1
        sq = ((rx ** 2) * (ry ** 2) - (rx ** 2) * (y1p ** 2) - (ry ** 2) * (x1p ** 2)) / \
             ((rx ** 2) * (y1p ** 2) + (ry ** 2) * (x1p ** 2))
        sq = max(0, sq)  # 防止数值误差
        coef = sign * math.sqrt(sq)

        cxp = coef * (rx * y1p / ry)
        cyp = coef * (-ry * x1p / rx)

        # 6. 转换回原坐标系
        cx = math.cos(phi) * cxp - math.sin(phi) * cyp + (x1 + x2) / 2
        cy = math.sin(phi) * cxp + math.cos(phi) * cyp + (y1 + y2) / 2

        # 7. 计算起始角度和角度差
        def vector_angle(ux, uy, vx, vy):
            dot = ux * vx + uy * vy
            mod = math.sqrt(ux * ux + uy * uy) * math.sqrt(vx * vx + vy * vy)
            if mod == 0:
                return 0
            angle = math.acos(max(-1, min(1, dot / mod)))
            if ux * vy - uy * vx < 0:
                angle = -angle
            return angle

        theta1 = vector_angle(1, 0, (x1p - cxp) / rx, (y1p - cyp) / ry)
        delta_theta = vector_angle((x1p - cxp) / rx, (y1p - cyp) / ry,
                                   (-x1p - cxp) / rx, (-y1p - cyp) / ry)

        if sweep == 0 and delta_theta > 0:
            delta_theta -= 2 * math.pi
        elif sweep == 1 and delta_theta < 0:
            delta_theta += 2 * math.pi

        # 8. 生成弧上的点
        points = []
        for i in range(steps + 1):
            t = i / steps
            theta = theta1 + delta_theta * t

            # 椭圆上的点（在椭圆坐标系）
            xe = rx * math.cos(theta)
            ye = ry * math.sin(theta)

            # 旋转并平移到原坐标系
            x = math.cos(phi) * xe - math.sin(phi) * ye + cx
            y = math.sin(phi) * xe + math.cos(phi) * ye + cy
            points.append((x, y))

        return points

    def _extract_path_points(self, d: str) -> List[Tuple[float, float]]:
        """从路径数据中提取所有坐标点 - 只提取第一个子路径"""
        import re

        points = []

        # 解析所有命令
        commands = re.findall(r'([MmLlHhVvCcSsQqTtAaZz])\s*([^MmLlHhVvCcSsQqTtAaZz]*)', d)

        current_pos = (0, 0)
        start_pos = (0, 0)

        for cmd, params_str in commands:
            params = self._parse_numbers(params_str)

            if cmd == 'M':  # 绝对移动到 - 如果是第二个M，停止提取
                if points:  # 已经有数据了，遇到新的M就停止
                    break
                if len(params) >= 2:
                    current_pos = (params[0], params[1])
                    start_pos = current_pos
                    points.append(current_pos)
                    # 后续参数视为L命令
                    for i in range(2, len(params), 2):
                        if i + 1 < len(params):
                            current_pos = (params[i], params[i + 1])
                            points.append(current_pos)

            elif cmd == 'm':  # 相对移动
                if points:
                    break
                if len(params) >= 2:
                    current_pos = (current_pos[0] + params[0], current_pos[1] + params[1])
                    start_pos = current_pos
                    points.append(current_pos)
                    for i in range(2, len(params), 2):
                        if i + 1 < len(params):
                            current_pos = (current_pos[0] + params[i], current_pos[1] + params[i + 1])
                            points.append(current_pos)

            elif cmd == 'L' or cmd == 'l':  # 直线到
                for i in range(0, len(params), 2):
                    if i + 1 < len(params):
                        if cmd == 'L':
                            current_pos = (params[i], params[i + 1])
                        else:
                            current_pos = (current_pos[0] + params[i], current_pos[1] + params[i + 1])
                        points.append(current_pos)

            elif cmd == 'H' or cmd == 'h':  # 水平线
                for x in params:
                    if cmd == 'H':
                        current_pos = (x, current_pos[1])
                    else:
                        current_pos = (current_pos[0] + x, current_pos[1])
                    points.append(current_pos)

            elif cmd == 'V' or cmd == 'v':  # 垂直线
                for y in params:
                    if cmd == 'V':
                        current_pos = (current_pos[0], y)
                    else:
                        current_pos = (current_pos[0], current_pos[1] + y)
                    points.append(current_pos)

            elif cmd == 'Z' or cmd == 'z':  # 闭合
                if start_pos != current_pos:
                    points.append(start_pos)
                current_pos = start_pos

        return points

    def _is_rect_path(self, d: str) -> bool:
        """检查是否为矩形路径"""
        import re
        # 检测矩形路径模式：M x y H x V y H x Z
        return bool(re.match(r'^[Mm][\d\s.,-]+[Hh][\d\s.,-]+[Vv][\d\s.,-]+[Hh][\d\s.,-]+[Zz]$', d))

    def _create_freeform_from_path(self, d: str, element: SVGElement, context: RenderContext):
        """从复杂路径创建自由形状"""
        points = self._extract_path_points(d)

        if len(points) < 2:
            return None

        from pptx.shapes.freeform import FreeformBuilder
        from pptx.util import Emu

        inch_points = [
            (
                context.offset_x + context.svg_to_inches(p[0]),
                context.offset_y + context.svg_to_inches(p[1])
            )
            for p in points
        ]

        # 计算包围盒和宽高
        min_x = min(p[0] for p in inch_points)
        min_y = min(p[1] for p in inch_points)
        max_x = max(p[0] for p in inch_points)
        max_y = max(p[1] for p in inch_points)

        width = max_x - min_x
        height = max_y - min_y

        # 避免零宽高
        if width < 0.001:
            width = 0.1
        if height < 0.001:
            height = 0.1

        # EMU 转换因子
        EMU_PER_INCH = 914400

        # 创建自由形状构建器
        builder = FreeformBuilder(
            context.slide.shapes,
            Emu(int(min_x * EMU_PER_INCH)),
            Emu(int(min_y * EMU_PER_INCH)),
            Emu(int(width * EMU_PER_INCH)),
            Emu(int(height * EMU_PER_INCH))
        )

        # 将点转换为相对于包围盒的归一化坐标
        segments = []
        for x, y in inch_points[1:]:
            norm_x = (x - min_x) / width if width > 0 else 0
            norm_y = (y - min_y) / height if height > 0 else 0
            segments.append((norm_x, norm_y))

        if segments:
            builder.add_line_segments(segments)

        shape = builder.convert_to_shape()

        style = context.parent_style.merge(element.style)
        if style.stroke and style.stroke.lower() != 'none':
            stroke_color = ColorParser.parse(style.stroke)
            if stroke_color:
                shape.line.color.rgb = stroke_color
        else:
            shape.line.color.rgb = RGBColor(0x7f, 0x8c, 0x8d)

        shape.line.width = Pt(ColorParser.parse_length(str(style.stroke_width)) or 1.5)

        # 填充处理 - 与基本形状一致
        if style.fill and style.fill.lower() == 'none':
            shape.fill.background()
        elif style.fill:
            fill_color = ColorParser.parse(style.fill)
            if fill_color:
                shape.fill.solid()
                shape.fill.fore_color.rgb = fill_color
        else:
            # 默认白色填充
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(255, 255, 255)

        return shape
