"""
主转换器
SVG到PPTX的完整转换流程
"""

import os
from pptx import Presentation
from pptx.util import Inches

from .parser import SVGParser
from .handlers import HandlerRegistry, RenderContext
from .handlers.basic_shapes import RectHandler, CircleHandler, EllipseHandler, LineHandler
from .handlers.polygons import PolygonHandler, PathHandler
from .handlers.text import TextHandler
from .handlers.group import GroupHandler


class SVGToPPTXConverter:
    """SVG到PPTX转换器 - 主类"""

    def __init__(self, svg_path: str):
        """
        初始化转换器

        Args:
            svg_path: SVG文件路径
        """
        self.svg_path = svg_path
        self.parser = SVGParser(svg_path)
        self.registry = self._init_registry()

    def _init_registry(self) -> HandlerRegistry:
        """初始化处理器注册表"""
        registry = HandlerRegistry()
        registry.clear()

        # 注册所有处理器（按优先级）
        registry.register(RectHandler(), priority=100)
        registry.register(CircleHandler(), priority=100)
        registry.register(EllipseHandler(), priority=100)
        registry.register(LineHandler(), priority=100)
        registry.register(PolygonHandler(), priority=90)
        registry.register(PathHandler(), priority=80)
        registry.register(TextHandler(), priority=100)
        registry.register(GroupHandler(), priority=50)

        return registry

    def convert(self, output_path: str = None,
                slide_width: float = 13.333,
                slide_height: float = 10.0,
                margin: float = 0.5) -> str:
        """
        执行SVG到PPTX的转换

        Args:
            output_path: 输出PPTX路径（默认与SVG同名）
            slide_width: 幻灯片宽度（英寸）
            slide_height: 幻灯片高度（英寸）
            margin: 边距（英寸）

        Returns:
            输出文件路径
        """
        # 确定输出路径
        if output_path is None:
            output_path = self.svg_path.replace('.svg', '.pptx')

        print(f"=" * 60)
        print(f"SVG to PPTX Converter v2.0")
        print(f"=" * 60)
        print(f"输入: {self.svg_path}")
        print(f"输出: {output_path}")
        print(f"=" * 60)

        # 解析SVG
        print(f"\n[1/3] 解析SVG文件...")
        elements = self.parser.parse()
        print(f"      找到 {len(elements)} 个顶层元素")
        print(f"      SVG尺寸: {self.parser.width:.1f} x {self.parser.height:.1f}")

        # 创建PPT
        print(f"\n[2/3] 创建PowerPoint...")
        prs = Presentation()
        prs.slide_width = Inches(slide_width)
        prs.slide_height = Inches(slide_height)

        blank_layout = prs.slide_layouts[6]  # 空白布局
        slide = prs.slides.add_slide(blank_layout)

        # 计算缩放和偏移
        available_width = slide_width - margin * 2
        available_height = slide_height - margin * 2

        svg_width_inch = self.parser.width / 96.0
        svg_height_inch = self.parser.height / 96.0

        scale_x = available_width / svg_width_inch
        scale_y = available_height / svg_height_inch
        scale = min(scale_x, scale_y)

        scaled_width = svg_width_inch * scale
        scaled_height = svg_height_inch * scale
        offset_x = (slide_width - scaled_width) / 2
        offset_y = (slide_height - scaled_height) / 2

        print(f"      缩放比例: {scale:.3f}")
        print(f"      幻灯片尺寸: {slide_width:.2f}\" x {slide_height:.2f}\"")

        # 创建渲染上下文
        context = RenderContext(
            slide=slide,
            scale=scale,
            offset_x=offset_x,
            offset_y=offset_y
        )

        # 处理每个元素
        print(f"\n[3/3] 转换元素...")
        success_count = 0
        error_count = 0

        for i, element in enumerate(elements, 1):
            handler = self.registry.get_handler(element)
            if handler:
                try:
                    result = handler.handle(element, context)
                    if result:
                        success_count += 1
                        if isinstance(result, list):
                            print(f"      [{i}/{len(elements)}] {element.tag}: +{len(result)}个子元素")
                        else:
                            print(f"      [{i}/{len(elements)}] {element.tag}: OK")
                    else:
                        print(f"      [{i}/{len(elements)}] {element.tag}: 跳过")
                except Exception as e:
                    error_count += 1
                    print(f"      [{i}/{len(elements)}] {element.tag}: 错误 - {e}")
            else:
                print(f"      [{i}/{len(elements)}] {element.tag}: 无处理器")

        # 保存
        prs.save(output_path)
        print(f"\n" + "=" * 60)
        print(f"转换完成!")
        print(f"  成功: {success_count}")
        print(f"  失败: {error_count}")
        print(f"  输出: {output_path}")
        print(f"=" * 60)

        return output_path


def convert_svg_to_pptx(svg_path: str, output_path: str = None, **kwargs) -> str:
    """
    便捷的转换函数

    Args:
        svg_path: SVG文件路径
        output_path: 输出PPTX路径（默认与SVG同名）
        **kwargs: 其他转换参数

    Returns:
        输出文件路径
    """
    converter = SVGToPPTXConverter(svg_path)
    return converter.convert(output_path, **kwargs)
