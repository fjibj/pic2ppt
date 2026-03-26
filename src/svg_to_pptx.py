"""
SVG to PowerPoint Converter
将SVG文件转换为PPTX中的可编辑形状
"""

import sys
import xml.etree.ElementTree as ET
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import nsmap
import re
from typing import Dict, List, Tuple, Optional
import os


class SVGElement:
    """SVG元素基类"""
    def __init__(self, tag: str, attrib: Dict[str, str], style: Dict[str, str]):
        self.tag = tag
        self.attrib = attrib
        self.style = style


class SVGRect(SVGElement):
    """SVG矩形"""
    def __init__(self, attrib: Dict[str, str], style: Dict[str, str]):
        super().__init__('rect', attrib, style)
        self.x = self._parse_length(attrib.get('x', '0'))
        self.y = self._parse_length(attrib.get('y', '0'))
        self.width = self._parse_length(attrib.get('width', '0'))
        self.height = self._parse_length(attrib.get('height', '0'))
        self.rx = self._parse_length(attrib.get('rx', '0'))
        # 解析虚线样式
        self.stroke_dasharray = style.get('stroke-dasharray', '') or attrib.get('stroke-dasharray', '')

    def _parse_length(self, value: str) -> float:
        """解析长度值（支持px单位）"""
        if not value:
            return 0.0
        match = re.match(r'([\d.]+)(?:px)?', str(value))
        return float(match.group(1)) if match else 0.0


class SVGText(SVGElement):
    """SVG文本"""
    def __init__(self, attrib: Dict[str, str], style: Dict[str, str], text: str):
        super().__init__('text', attrib, style)
        self.x = self._parse_length(attrib.get('x', '0'))
        self.y = self._parse_length(attrib.get('y', '0'))
        self.text_anchor = attrib.get('text-anchor', 'start')
        self.dominant_baseline = attrib.get('dominant-baseline', 'auto')
        self.text = text
        self.font_size = self._parse_font_size(style.get('font-size', '16px'))
        self.font_weight = style.get('font-weight', '400')

    def _parse_length(self, value: str) -> float:
        if not value:
            return 0.0
        match = re.match(r'([\d.]+)(?:px)?', str(value))
        return float(match.group(1)) if match else 0.0

    def _parse_font_size(self, value: str) -> float:
        match = re.match(r'([\d.]+)px', str(value))
        return float(match.group(1)) if match else 16.0


class SVGLine(SVGElement):
    """SVG线条"""
    def __init__(self, attrib: Dict[str, str], style: Dict[str, str]):
        super().__init__('line', attrib, style)
        self.x1 = self._parse_length(attrib.get('x1', '0'))
        self.y1 = self._parse_length(attrib.get('y1', '0'))
        self.x2 = self._parse_length(attrib.get('x2', '0'))
        self.y2 = self._parse_length(attrib.get('y2', '0'))
        # 检测是否有箭头标记
        marker_end = attrib.get('marker-end', '')
        self.has_arrow = 'arrow' in marker_end.lower() or 'url(#arrow)' in marker_end

    def _parse_length(self, value: str) -> float:
        if not value:
            return 0.0
        match = re.match(r'([\d.]+)(?:px)?', str(value))
        return float(match.group(1)) if match else 0.0


class SVGPolygon(SVGElement):
    """SVG多边形（包括菱形/钻石形）"""
    def __init__(self, attrib: Dict[str, str], style: Dict[str, str]):
        super().__init__('polygon', attrib, style)
        self.points = self._parse_points(attrib.get('points', ''))
        self.is_diamond = self._is_diamond()
        # 计算边界框
        if self.points:
            xs = [p[0] for p in self.points]
            ys = [p[1] for p in self.points]
            self.x = min(xs)
            self.y = min(ys)
            self.width = max(xs) - min(xs)
            self.height = max(ys) - min(ys)
            self.center_x = (min(xs) + max(xs)) / 2
            self.center_y = (min(ys) + max(ys)) / 2
        else:
            self.x = self.y = self.width = self.height = 0
            self.center_x = self.center_y = 0

    def _parse_points(self, points_str: str) -> List[Tuple[float, float]]:
        """解析points属性，格式: "x1,y1 x2,y2 x3,y3 ..."""
        points = []
        if not points_str:
            return points
        # 替换逗号为空格，然后分割
        coords = points_str.replace(',', ' ').split()
        for i in range(0, len(coords) - 1, 2):
            try:
                x = float(coords[i])
                y = float(coords[i + 1])
                points.append((x, y))
            except (ValueError, IndexError):
                continue
        return points

    def _is_diamond(self) -> bool:
        """判断是否为菱形（4个点，对角线对齐）"""
        if len(self.points) != 4:
            return False
        # 菱形的特征：4个点，对角线交点为中心
        return True

    def _parse_length(self, value: str) -> float:
        if not value:
            return 0.0
        match = re.match(r'([\d.]+)(?:px)?', str(value))
        return float(match.group(1)) if match else 0.0


class SVGParser:
    """SVG解析器"""

    def __init__(self, svg_path: str):
        self.svg_path = svg_path
        self.elements: List[SVGElement] = []
        self.viewbox: Tuple[float, float, float, float] = (0, 0, 680, 860)
        self.width = 680
        self.height = 860

    def parse(self) -> List[SVGElement]:
        """解析SVG文件"""
        tree = ET.parse(self.svg_path)
        root = tree.getroot()

        # 获取viewBox和尺寸
        viewbox_attr = root.get('viewBox')
        if viewbox_attr:
            parts = viewbox_attr.split()
            self.viewbox = tuple(float(p) for p in parts)
            self.width = self.viewbox[2]
            self.height = self.viewbox[3]

        # 解析所有元素
        self._parse_element(root)

        return self.elements

    def _parse_style(self, attrib: Dict[str, str]) -> Dict[str, str]:
        """解析内联样式"""
        style = {}

        # 直接属性
        for key in ['fill', 'stroke', 'stroke-width', 'font-size', 'font-weight',
                   'opacity', 'font-family', 'color']:
            if key in attrib:
                style[key] = attrib[key]

        # style属性
        style_attr = attrib.get('style', '')
        if style_attr:
            for item in style_attr.split(';'):
                if ':' in item:
                    key, value = item.split(':', 1)
                    style[key.strip()] = value.strip()

        return style

    def _parse_element(self, element: ET.Element):
        """递归解析SVG元素"""
        tag = element.tag.split('}')[-1] if '}' in element.tag else element.tag

        if tag == 'rect':
            style = self._parse_style(element.attrib)
            self.elements.append(SVGRect(element.attrib, style))
        elif tag == 'text':
            style = self._parse_style(element.attrib)
            text = ''.join(element.itertext())
            if text.strip():
                self.elements.append(SVGText(element.attrib, style, text.strip()))
        elif tag == 'line':
            style = self._parse_style(element.attrib)
            self.elements.append(SVGLine(element.attrib, style))
        elif tag == 'polygon':
            style = self._parse_style(element.attrib)
            self.elements.append(SVGPolygon(element.attrib, style))

        # 递归解析子元素
        for child in element:
            self._parse_element(child)


class SVGToPPTXConverter:
    """SVG转PPTX转换器"""

    def __init__(self, svg_path: str):
        self.svg_path = svg_path
        self.parser = SVGParser(svg_path)
        self.elements = []
        # 缩放比例：将SVG像素转换为PPT英寸 (96 DPI)
        self.scale = 914400 / 96 / 914400  # 1 px = 1/96 inch

    def _svg_to_inches(self, value: float) -> float:
        """将SVG像素值转换为英寸"""
        return value / 96.0  # 标准96 DPI

    def _parse_color(self, color_str: str) -> Optional[RGBColor]:
        """解析颜色字符串为RGB"""
        if not color_str or color_str == 'none':
            return None

        # rgba(r, g, b, a)格式 - 忽略透明度
        match = re.match(r'rgba\((\d+),\s*(\d+),\s*(\d+),\s*([\d.]+)\)', color_str)
        if match:
            return RGBColor(int(match.group(1)), int(match.group(2)), int(match.group(3)))

        # rgb(r, g, b)格式
        match = re.match(r'rgb\((\d+),\s*(\d+),\s*(\d+)\)', color_str)
        if match:
            return RGBColor(int(match.group(1)), int(match.group(2)), int(match.group(3)))

        # 十六进制格式
        if color_str.startswith('#'):
            hex_color = color_str[1:]
            if len(hex_color) == 6:
                r = int(hex_color[0:2], 16)
                g = int(hex_color[2:4], 16)
                b = int(hex_color[4:6], 16)
                return RGBColor(r, g, b)
            elif len(hex_color) == 3:
                r = int(hex_color[0]*2, 16)
                g = int(hex_color[1]*2, 16)
                b = int(hex_color[2]*2, 16)
                return RGBColor(r, g, b)

        # 命名颜色
        color_map = {
            'black': RGBColor(0, 0, 0),
            'white': RGBColor(255, 255, 255),
            'red': RGBColor(255, 0, 0),
            'green': RGBColor(0, 128, 0),
            'blue': RGBColor(0, 0, 255),
            'yellow': RGBColor(255, 255, 0),
            'gray': RGBColor(128, 128, 128),
        }
        return color_map.get(color_str.lower())

    def _extract_stroke_width(self, width_str: str) -> float:
        """解析线宽"""
        if not width_str:
            return 1.0
        match = re.match(r'([\d.]+)(?:px)?', str(width_str))
        return float(match.group(1)) if match else 1.0

    def convert(self, output_path: str, slide_width: float = 13.333, slide_height: float = 10):
        """
        转换SVG为PPTX

        Args:
            output_path: 输出PPTX文件路径
            slide_width: 幻灯片宽度（英寸），默认13.333英寸(16:9)
            slide_height: 幻灯片高度（英寸），默认10英寸
        """
        # 解析SVG
        print(f"正在解析SVG文件: {self.svg_path}")
        self.elements = self.parser.parse()
        print(f"解析完成，共 {len(self.elements)} 个元素")

        # 创建演示文稿
        prs = Presentation()
        prs.slide_width = Inches(slide_width)
        prs.slide_height = Inches(slide_height)

        # 添加空白幻灯片
        blank_layout = prs.slide_layouts[6]  # 空白布局
        slide = prs.slides.add_slide(blank_layout)

        # 计算缩放以适应幻灯片
        svg_width = self.parser.width
        svg_height = self.parser.height

        # 居中放置
        scale_x = (slide_width - 1) / self._svg_to_inches(svg_width)
        scale_y = (slide_height - 1) / self._svg_to_inches(svg_height)
        scale = min(scale_x, scale_y)

        # 居中偏移
        offset_x = (slide_width - self._svg_to_inches(svg_width) * scale) / 2
        offset_y = (slide_height - self._svg_to_inches(svg_height) * scale) / 2

        print(f"缩放比例: {scale:.3f}, 偏移: ({offset_x:.2f}, {offset_y:.2f})")

        # 转换每个元素
        for elem in self.elements:
            try:
                self._convert_element(slide, elem, scale, offset_x, offset_y)
            except Exception as e:
                print(f"转换元素 {elem.tag} 时出错: {e}")

        # 保存
        prs.save(output_path)
        print(f"PPTX已保存: {output_path}")

    def _convert_element(self, slide, elem: SVGElement, scale: float, offset_x: float, offset_y: float):
        """转换单个元素"""
        if isinstance(elem, SVGRect):
            self._convert_rect(slide, elem, scale, offset_x, offset_y)
        elif isinstance(elem, SVGText):
            self._convert_text(slide, elem, scale, offset_x, offset_y)
        elif isinstance(elem, SVGLine):
            self._convert_line(slide, elem, scale, offset_x, offset_y)
        elif isinstance(elem, SVGPolygon):
            self._convert_polygon(slide, elem, scale, offset_x, offset_y)

    def _convert_rect(self, slide, rect: SVGRect, scale: float, offset_x: float, offset_y: float):
        """转换矩形"""
        left = Inches(offset_x + self._svg_to_inches(rect.x) * scale)
        top = Inches(offset_y + self._svg_to_inches(rect.y) * scale)
        width = Inches(self._svg_to_inches(rect.width) * scale)
        height = Inches(self._svg_to_inches(rect.height) * scale)

        # 选择形状类型
        if rect.rx > 0:
            # 圆角矩形
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left, top, width, height
            )
        else:
            # 普通矩形
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left, top, width, height
            )

        # 设置填充颜色
        fill_color = self._parse_color(rect.style.get('fill'))
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
        else:
            shape.fill.background()

        # 设置边框（包括虚线）
        stroke_color = self._parse_color(rect.style.get('stroke'))
        stroke_width = self._extract_stroke_width(rect.style.get('stroke-width', '1'))
        has_dash = rect.stroke_dasharray and rect.stroke_dasharray != 'none'

        if stroke_color or has_dash:
            if stroke_color:
                shape.line.color.rgb = stroke_color
            else:
                shape.line.color.rgb = RGBColor(128, 128, 128)  # 默认灰色
            shape.line.width = Pt(stroke_width)

            # 设置虚线样式
            if has_dash:
                shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        else:
            shape.line.fill.background()

    def _convert_polygon(self, slide, poly: SVGPolygon, scale: float, offset_x: float, offset_y: float):
        """转换多边形（菱形等）"""
        left = Inches(offset_x + self._svg_to_inches(poly.x) * scale)
        top = Inches(offset_y + self._svg_to_inches(poly.y) * scale)
        width = Inches(self._svg_to_inches(poly.width) * scale)
        height = Inches(self._svg_to_inches(poly.height) * scale)

        # 根据类型选择形状
        if poly.is_diamond:
            # 菱形
            shape = slide.shapes.add_shape(
                MSO_SHAPE.DIAMOND,
                left, top, width, height
            )
        else:
            # 其他多边形使用矩形代替
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left, top, width, height
            )

        # 设置填充颜色
        fill_color = self._parse_color(poly.style.get('fill'))
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
        else:
            shape.fill.background()

        # 设置边框
        stroke_color = self._parse_color(poly.style.get('stroke'))
        if stroke_color:
            shape.line.color.rgb = stroke_color
            shape.line.width = Pt(self._extract_stroke_width(poly.style.get('stroke-width', '1')))
        else:
            shape.line.fill.background()

    def _convert_text(self, slide, text: SVGText, scale: float, offset_x: float, offset_y: float):
        """转换文本"""
        # 字体大小（pt）
        font_size_pt = max(text.font_size * scale * 0.75, 8)  # px到pt的近似转换，最小8pt

        # 估算文本尺寸（英寸）
        # 中文字符宽度约为字体大小的1.1倍（以pt为单位），再转换为英寸（/72）
        char_width_inch = (font_size_pt * 1.1) / 72
        text_width_inch = len(text.text) * char_width_inch
        # 文本高度 = 字体高度 + 行间距
        text_height_inch = (font_size_pt * 1.3) / 72

        # 计算基础位置（英寸）
        base_x = offset_x + self._svg_to_inches(text.x) * scale
        base_y = offset_y + self._svg_to_inches(text.y) * scale

        # 根据text-anchor调整水平位置
        if text.text_anchor == 'middle':
            left = Inches(base_x - text_width_inch / 2)
        elif text.text_anchor == 'end':
            left = Inches(base_x - text_width_inch)
        else:  # start
            left = Inches(base_x)

        # 垂直位置：SVG的y是基线位置，PPT的top是文本框顶部
        # 需要向上偏移约0.8倍字体高度来对齐基线
        baseline_offset = (font_size_pt * 0.8) / 72
        top = Inches(base_y - baseline_offset)

        # 添加文本框
        txBox = slide.shapes.add_textbox(
            left, top,
            Inches(max(text_width_inch, 0.05)),
            Inches(max(text_height_inch, 0.05))
        )

        tf = txBox.text_frame
        tf.text = text.text

        # 设置字体
        p = tf.paragraphs[0]
        p.font.size = Pt(font_size_pt)
        p.font.name = 'Microsoft YaHei'  # 使用微软雅黑支持中文

        # 设置颜色
        fill_color = self._parse_color(text.style.get('fill'))
        if fill_color:
            p.font.color.rgb = fill_color

        # 设置对齐
        if text.text_anchor == 'middle':
            p.alignment = PP_ALIGN.CENTER
        elif text.text_anchor == 'end':
            p.alignment = PP_ALIGN.RIGHT
        else:
            p.alignment = PP_ALIGN.LEFT

    def _convert_line(self, slide, line: SVGLine, scale: float, offset_x: float, offset_y: float):
        """转换线条"""
        begin_x = Inches(offset_x + self._svg_to_inches(line.x1) * scale)
        begin_y = Inches(offset_y + self._svg_to_inches(line.y1) * scale)
        end_x = Inches(offset_x + self._svg_to_inches(line.x2) * scale)
        end_y = Inches(offset_y + self._svg_to_inches(line.y2) * scale)

        # 添加线条
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            begin_x, begin_y, end_x, end_y
        )

        # 设置线条样式
        stroke_color = self._parse_color(line.style.get('stroke'))
        if stroke_color:
            connector.line.color.rgb = stroke_color
        else:
            connector.line.color.rgb = RGBColor(115, 114, 108)  # 默认灰色

        connector.line.width = Pt(self._extract_stroke_width(line.style.get('stroke-width', '1.5')))

        # 添加箭头头部（如果标记了arrow）
        if line.has_arrow:
            self._add_arrow_head(slide, line, scale, offset_x, offset_y, stroke_color or RGBColor(115, 114, 108))

    def _add_arrow_head(self, slide, line: SVGLine, scale: float, offset_x: float, offset_y: float, color: RGBColor):
        """在线条终点添加箭头头部"""
        import math

        # 计算箭头位置（线条终点）
        end_x = offset_x + self._svg_to_inches(line.x2) * scale
        end_y = offset_y + self._svg_to_inches(line.y2) * scale

        # 计算线条角度
        dx = line.x2 - line.x1
        dy = line.y2 - line.y1
        angle = math.atan2(dy, dx)

        # 箭头大小
        arrow_size = 0.12 * scale  # 英寸

        # 添加三角形箭头
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.ISOSCELES_TRIANGLE,
            Inches(end_x - arrow_size/2),
            Inches(end_y - arrow_size/2),
            Inches(arrow_size),
            Inches(arrow_size)
        )

        # 旋转箭头以指向正确的方向
        # 三角形默认顶点朝上，需要旋转使顶点指向线条方向
        arrow.rotation = math.degrees(angle) + 90

        # 设置箭头样式
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = color
        arrow.line.fill.background()  # 无边框

        # 将箭头移到线条下方（可选，取决于视觉效果）


def main():
    """主函数"""
    # 获取脚本所在目录
    script_dir = os.path.dirname(os.path.abspath(__file__))
    # 项目根目录（脚本目录的父目录）
    project_dir = os.path.dirname(script_dir)

    svg_file = os.path.join(project_dir, "ai_component_management_platform.svg")
    output_file = os.path.join(project_dir, "output_architecture.pptx")

    if not os.path.exists(svg_file):
        print("错误: SVG文件不存在: {}".format(svg_file))
        print("请确保文件存在于项目根目录")
        return

    print("输入SVG: {}".format(svg_file))
    print("输出PPTX: {}".format(output_file))

    try:
        converter = SVGToPPTXConverter(svg_file)
        converter.convert(output_file)
        print("转换完成！")
    except Exception as e:
        print("转换失败: {}".format(e))
        import traceback
        traceback.print_exc()


if __name__ == "__main__":
    main()
